import os
import io
from lxml import etree
from pptx import Presentation as PptxPresentation
from pptx.oxml.ns import qn
from PIL import Image, ImageDraw, ImageFont
from ppt_render_engine.core.pptx_builder import build_pptx
from ppt_render_engine.models.schema import Presentation
from ppt_render_engine.log_config import app_logger
from ppt_render_engine.temp_util import get_temp_path

logger = app_logger("engine.render")


def _load_font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        r"C:\Windows\Fonts\msyh.ttc",
        r"C:\Windows\Fonts\msyhbd.ttc",
        r"C:\Windows\Fonts\simsun.ttc",
        r"C:\Windows\Fonts\arial.ttf",
        r"C:\Windows\Fonts\segoeui.ttf",
    ]
    for p in candidates:
        if os.path.isfile(p):
            try:
                return ImageFont.truetype(p, size)
            except Exception:
                continue
    return ImageFont.load_default()


_FONT_CACHE: dict[int, ImageFont.FreeTypeFont | ImageFont.ImageFont] = {}


def _get_font(size: int):
    if size not in _FONT_CACHE:
        _FONT_CACHE[size] = _load_font(size)
    return _FONT_CACHE[size]


def _parse_color(val: str | None) -> tuple[int, int, int] | None:
    if not val:
        return None
    try:
        hex_str = val.lstrip("#")
        return (int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
    except Exception:
        return None


def _get_theme_colors(prs: PptxPresentation) -> dict[str, str]:
    try:
        sm_part = prs.slide_masters[0].part
        for rel in sm_part.rels.values():
            if rel.reltype == "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme":
                theme = etree.fromstring(rel.target_part.blob)
                clr = theme.find('.//' + qn('a:clrScheme'))
                if clr is None:
                    return {}
                colors = {}
                for child in clr:
                    tag = child.tag.split("}")[-1]
                    srgb = child.find(qn('a:srgbClr'))
                    if srgb is not None:
                        colors[tag] = srgb.get("val", "").upper()
                return colors
    except Exception:
        pass
    return {}


def _resolve_fill(shape_elem, theme_colors: dict) -> tuple[int, int, int] | None:
    spPr = shape_elem.find(qn("p:spPr"))
    if spPr is None:
        return None
    solid = spPr.find(qn("a:solidFill"))
    if solid is None:
        return None
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is not None:
        return _parse_color(srgb.get("val"))
    scheme = solid.find(qn("a:schemeClr"))
    if scheme is not None and scheme.get("val") in theme_colors:
        return _parse_color(theme_colors[scheme.get("val")])
    return None


def _resolve_line(shape_elem, theme_colors: dict) -> tuple[int, int, int] | None:
    spPr = shape_elem.find(qn("p:spPr"))
    if spPr is None:
        return None
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        return None
    solid = ln.find(qn("a:solidFill"))
    if solid is None:
        return None
    srgb = solid.find(qn("a:srgbClr"))
    if srgb is not None:
        return _parse_color(srgb.get("val"))
    scheme = solid.find(qn("a:schemeClr"))
    if scheme is not None and scheme.get("val") in theme_colors:
        return _parse_color(theme_colors[scheme.get("val")])
    return None


def _has_no_fill(shape_elem) -> bool:
    spPr = shape_elem.find(qn("p:spPr"))
    if spPr is None:
        return False
    return spPr.find(qn("a:noFill")) is not None


def _emus_to_px(emu: int, slide_dim: int, img_dim: int) -> int:
    return int(emu / slide_dim * img_dim) if slide_dim else 0


def _pt_to_px(pt: float, dpi: int) -> int:
    return max(6, int(pt * dpi / 72))


def render_preview(presentation: Presentation, slide_index: int = 0, dpi: int = 150) -> bytes:
    tmp_path = get_temp_path(".pptx")
    logger.info("Rendering preview via Pillow", slide=slide_index, dpi=dpi)
    try:
        build_pptx(presentation, tmp_path)
        prs = PptxPresentation(tmp_path)
        slide_w = int(prs.slide_width)
        slide_h = int(prs.slide_height)
        img_w = max(1, int(slide_w / 914400 * dpi))
        img_h = max(1, int(slide_h / 914400 * dpi))
        theme_colors = _get_theme_colors(prs)

        img = Image.new("RGB", (img_w, img_h), (255, 255, 255))
        draw = ImageDraw.Draw(img)

        bg_elem = prs.slides[slide_index].background._element
        bg_fill = _resolve_fill(bg_elem, theme_colors)
        if bg_fill:
            img = Image.new("RGB", (img_w, img_h), bg_fill)
            draw = ImageDraw.Draw(img)

        pptx_slide = prs.slides[slide_index]
        for s in pptx_slide.shapes:
            l = _emus_to_px(s.left, slide_w, img_w) if s.left else 0
            t = _emus_to_px(s.top, slide_h, img_h) if s.top else 0
            w = _emus_to_px(s.width, slide_w, img_w) if s.width else 0
            h = _emus_to_px(s.height, slide_h, img_h) if s.height else 0
            if w < 2 or h < 2:
                continue

            fill_c = _resolve_fill(s._element, theme_colors)
            border_c = _resolve_line(s._element, theme_colors)

            if fill_c:
                draw.rounded_rectangle([l, t, l + w, t + h], radius=2, fill=fill_c)
            if border_c:
                draw.rounded_rectangle([l, t, l + w, t + h], radius=2, fill=None,
                                       outline=border_c, width=max(1, int(w * 0.005)))

            # ----- Tables -----
            if s.shape_type == 13:
                try:
                    tbl = s.table
                    tbl_elem = tbl._tbl
                    grid = tbl_elem.find(qn('a:tblGrid'))
                    col_widths = []
                    if grid is not None:
                        for col in grid.findall(qn('a:gridCol')):
                            cw = col.get('w')
                            col_widths.append(int(cw) if cw else 0)

                    row_heights = []
                    for row in tbl.rows:
                        row_heights.append(int(row.height or 0))

                    cy = t
                    for ri, row in enumerate(tbl.rows):
                        cx = l
                        rh = row_heights[ri] if ri < len(row_heights) else 0
                        for ci, cell in enumerate(row.cells):
                            cw = col_widths[ci] if ci < len(col_widths) else 0
                            if cw > 0 and rh > 0:
                                cell_x = _emus_to_px(cx, slide_w, img_w)
                                cell_y = _emus_to_px(cy, slide_h, img_h)
                                cell_w = _emus_to_px(cw, slide_w, img_w)
                                cell_h = _emus_to_px(rh, slide_h, img_h)
                                cell_fill = _resolve_fill(cell._tc, theme_colors)
                                if cell_fill:
                                    draw.rectangle([cell_x, cell_y, cell_x + cell_w, cell_y + cell_h], fill=cell_fill)
                                draw.rectangle([cell_x, cell_y, cell_x + cell_w, cell_y + cell_h],
                                               outline=(180, 180, 180), width=1)
                                txt = cell.text.strip()
                                if txt:
                                    draw.text((cell_x + 2, cell_y + 2), txt, fill=(30, 30, 30), font=_get_font(10))
                            cx += cw
                        cy += rh
                    continue
                except Exception as e:
                    logger.debug("Table render error", error=str(e))

            # ----- Text / shape -----
            if s.has_text_frame:
                tf = s.text_frame
                y_offset = t + 3
                for pi, para in enumerate(tf.paragraphs):
                    txt = "".join(r.text for r in para.runs)
                    if not txt.strip():
                        continue
                    try:
                        first_run = next((r for r in para.runs if r.text.strip()), None)
                    except StopIteration:
                        first_run = None
                    pt_size = 12
                    if first_run:
                        try:
                            fs = first_run.font.size
                            if fs:
                                pt_size = fs / 12700
                        except Exception:
                            pass
                    px_size = _pt_to_px(pt_size, dpi)
                    font = _get_font(px_size)
                    text_color = (30, 30, 30)
                    if first_run:
                        try:
                            rPr = first_run._r.find(qn("a:rPr"))
                            if rPr is not None:
                                sf = rPr.find(qn("a:solidFill"))
                                if sf is not None:
                                    srgb = sf.find(qn("a:srgbClr"))
                                    if srgb is not None:
                                        text_color = _parse_color(srgb.get("val")) or text_color
                                    scheme = sf.find(qn("a:schemeClr"))
                                    if scheme is not None and scheme.get("val") in theme_colors:
                                        text_color = _parse_color(theme_colors[scheme.get("val")]) or text_color
                        except Exception:
                            pass
                    draw.text((l + 3, y_offset), txt, fill=text_color, font=font)
                    y_offset += px_size + 3

        buf = io.BytesIO()
        img.save(buf, format="PNG")
        png_bytes = buf.getvalue()
        logger.debug("Preview rendered via Pillow", size_kb=round(len(png_bytes) / 1024, 1))
        return png_bytes
    finally:
        if os.path.exists(tmp_path):
            os.unlink(tmp_path)
