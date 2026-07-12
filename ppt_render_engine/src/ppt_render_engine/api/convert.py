import os
import json
import zipfile
import io
import uuid
from urllib.parse import quote
from lxml import etree
from fastapi import APIRouter, UploadFile, File, Query, HTTPException
from fastapi.responses import Response
from pptx import Presentation as PptxPresentation
from pptx.oxml.ns import qn
from ppt_render_engine.core.pptx_parser import parse_pptx
from ppt_render_engine.core.pptx_builder import build_pptx
from ppt_render_engine.core.template import get_template_manager
from ppt_render_engine.models.schema import Presentation
from ppt_render_engine.log_config import app_logger
from ppt_render_engine.temp_util import get_temp_path

router = APIRouter(prefix="/convert", tags=["Convert"])
logger = app_logger("api.convert")


def _strip_slides(src_path: str, dst_path: str) -> None:
    """Remove all content slides from a PPTX, keeping only layouts + theme."""
    prs = PptxPresentation(src_path)
    sldIdLst = prs.slides._sldIdLst
    for sld in list(sldIdLst):
        rId = sld.get(qn("r:id"))
        if rId:
            prs.part.drop_rel(rId)
        sldIdLst.remove(sld)
    prs.save(dst_path)


@router.post("/pptx-to-json")
async def pptx_to_json(file: UploadFile = File(...)):
    if not file.filename or not file.filename.endswith((".pptx", ".PPTX")):
        raise HTTPException(400, "Only .pptx files are supported")
    tmp_path = get_temp_path(".pptx")
    logger.info("Parsing PPTX to JSON", file=file.filename)
    try:
        raw = await file.read()
        with open(tmp_path, "wb") as f:
            f.write(raw)
        pres = parse_pptx(tmp_path)
        slide_count = len(pres.slides)
        # Extract color scheme name from source PPTX theme
        try:
            from pptx.opc.constants import RELATIONSHIP_TYPE as RT
            src_prs = PptxPresentation(tmp_path)
            if src_prs.slide_masters:
                sm_part = src_prs.slide_masters[0].part
                for rel in sm_part.rels.values():
                    if rel.reltype == RT.THEME:
                        theme = etree.fromstring(rel.target_part.blob)
                        clrScheme = theme.find('.//' + qn('a:clrScheme'))
                        if clrScheme is not None and clrScheme.get('name'):
                            pres.color_scheme = clrScheme.get('name')
                        break
        except Exception:
            pass
        logger.info("PPTX parsed successfully", file=file.filename, slides=slide_count, scheme=pres.color_scheme)
        json_bytes = pres.model_dump_json(exclude_none=True, indent=2).encode("utf-8")

        # Build zip: presentation.json + images/ + template.pptx (layouts only)
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
            zf.writestr("presentation.json", json_bytes)
            stripped = get_temp_path(".pptx")
            _strip_slides(tmp_path, stripped)
            with open(stripped, "rb") as f:
                zf.writestr("template.pptx", f.read())
            os.unlink(stripped)
            try:
                with zipfile.ZipFile(tmp_path, "r") as src:
                    for name in src.namelist():
                        if name.startswith("ppt/media/"):
                            zf.writestr(f"images/{os.path.basename(name)}", src.read(name))
            except Exception:
                pass
        buf.seek(0)
        stem = os.path.splitext(file.filename)[0]
        encoded = quote(f"{stem}.zip", safe="")
        return Response(
            content=buf.getvalue(),
            media_type="application/zip",
            headers={"Content-Disposition": f"attachment; filename*=UTF-8''{encoded}"},
        )
    except Exception as e:
        logger.error("PPTX parse failed", file=file.filename, error=str(e))
        raise HTTPException(500, detail="PPTX 解析失败，请确认文件是有效的 PowerPoint 格式")
    finally:
        os.unlink(tmp_path)


@router.post("/json-to-pptx")
async def json_to_pptx(
    file: UploadFile = File(...),
    template_name: str = Query(None, description="Override template_name from JSON body"),
    color_scheme: str = Query(None, description="Override color_scheme from JSON body"),
):
    if not file.filename or not file.filename.endswith((".json", ".JSON", ".zip", ".ZIP")):
        raise HTTPException(400, "Only .json or .zip files are supported")
    raw = await file.read()
    is_zip = file.filename.lower().endswith(".zip")

    template_cleanup = None
    if is_zip:
        image_dir = get_temp_path("_images")
        os.makedirs(image_dir, exist_ok=True)
        with zipfile.ZipFile(io.BytesIO(raw)) as zf:
            json_raw = zf.read("presentation.json")
            # Extract embedded template if present
            if "template.pptx" in zf.namelist():
                tm = get_template_manager()
                template_cleanup = f"_zip_{uuid.uuid4().hex}"
                tm.save_to_storage(template_cleanup, zf.read("template.pptx"))
                tm.load_from_storage(template_cleanup)
                logger.info("Embedded template loaded from zip", name=template_cleanup)
            for name in zf.namelist():
                if name.startswith("images/") and not name.endswith("/"):
                    dest = os.path.join(image_dir, os.path.basename(name))
                    with open(dest, "wb") as f:
                        f.write(zf.read(name))
        pres_dict = json.loads(json_raw)
        if template_cleanup:
            pres_dict["template_name"] = template_cleanup
        # Rewrite image paths to point at extracted files
        for slide in pres_dict.get("slides", []):
            for shape in slide.get("shapes", []):
                img = shape.get("image_content")
                if img and img.get("path"):
                    fname = os.path.basename(img["path"])
                    candidate = os.path.join(image_dir, fname)
                    if os.path.exists(candidate):
                        img["path"] = candidate
                        img["data"] = None
                        img["url"] = None
    else:
        pres_dict = json.loads(raw)

    pres = Presentation(**pres_dict)
    if template_name is not None:
        pres.template_name = template_name
        logger.info("Template overridden by query param", template=template_name)
    if color_scheme is not None:
        pres.color_scheme = color_scheme
        logger.info("Color scheme overridden by query param", scheme=color_scheme)
    tmp_path = get_temp_path(".pptx")
    slide_count = len(pres.slides)
    logger.info("Building PPTX from JSON file", file=file.filename, slides=slide_count, template=pres.template_name)
    try:
        build_pptx(pres, tmp_path)
        logger.info("PPTX built successfully", slides=slide_count)
        with open(tmp_path, "rb") as f:
            content = f.read()
        out_name = file.filename.replace(".zip", ".pptx").replace(".ZIP", ".pptx").replace(".json", ".pptx").replace(".JSON", ".pptx")
        encoded = quote(out_name, safe="")
        return Response(
            content=content,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename*=UTF-8''{encoded}"},
        )
    except Exception as e:
        logger.error("PPTX build failed", error=str(e))
        raise HTTPException(500, detail="PPTX 生成失败，请检查 JSON 中字段是否符合规范（如坐标不越界、字体存在等）")
    finally:
        if template_cleanup:
            tm = get_template_manager()
            tm.unload(template_cleanup)
            tm.delete_from_storage(template_cleanup)
            logger.info("Embedded template cleaned up", name=template_cleanup)



