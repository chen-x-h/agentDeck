import json
import os
from fastapi import APIRouter, HTTPException, Query, UploadFile, File, Request
from fastapi.responses import Response
from pptx import Presentation as PptxPresentation
from ppt_render_engine.models.schema import Presentation
from ppt_render_engine.core.pptx_builder import build_pptx
from ppt_render_engine.core.pptx_parser import _parse_shape, _parse_xml_fill
from ppt_render_engine.engine.render import render_preview
from ppt_render_engine.engine.layout_preset import PRESETS
from ppt_render_engine.log_config import app_logger
from ppt_render_engine.temp_util import get_temp_path

router = APIRouter(prefix="/render", tags=["Preview"])
logger = app_logger("api.preview")


@router.post("/preview")
async def preview_slide(
    file: UploadFile = File(...),
    slide_index: int = Query(0, ge=0),
    dpi: int = Query(150, ge=72, le=300),
):
    if not file.filename or not file.filename.lower().endswith((".json", ".zip")):
        raise HTTPException(400, detail="文件格式错误：仅支持 .json 或 .zip 文件")
    raw = await file.read()
    is_zip = file.filename.lower().endswith(".zip")
    if is_zip:
        import zipfile, io
        try:
            with zipfile.ZipFile(io.BytesIO(raw)) as zf:
                raw = zf.read("presentation.json")
        except Exception:
            raise HTTPException(400, detail="ZIP 解析失败：未找到 presentation.json")
    try:
        pres_dict = json.loads(raw)
    except json.JSONDecodeError:
        raise HTTPException(400, detail="JSON 解析失败：文件内容不是有效的 JSON 格式")
    data = Presentation(**pres_dict)
    if slide_index >= len(data.slides):
        raise HTTPException(400, detail=f"页码越界：共 {len(data.slides)} 页，请求第 {slide_index} 页")
    logger.info("Rendering preview from file", file=file.filename, slide=slide_index, dpi=dpi)
    try:
        png_bytes = render_preview(data, slide_index=slide_index, dpi=dpi)
        size_kb = len(png_bytes) / 1024
        logger.info("Preview rendered", slide=slide_index, size_kb=round(size_kb, 1))
        return Response(content=png_bytes, media_type="image/png")
    except Exception as e:
        logger.error("Preview failed", slide=slide_index, error=str(e))
        raise HTTPException(500, detail="预览图生成失败：" + str(e))


@router.get("/presets")
async def list_presets():
    result = []
    for name, zones in PRESETS.items():
        result.append({
            "name": name,
            "zones": [{"role": z.role, "left_pct": z.left_pct, "top_pct": z.top_pct,
                        "width_pct": z.width_pct, "height_pct": z.height_pct} for z in zones],
        })
    return {"presets": result}


@router.post("/preview-info")
async def preview_info(
    file: UploadFile = File(...),
    slide_index: int = Query(0, ge=0),
):
    if not file.filename or not file.filename.lower().endswith((".json", ".zip")):
        raise HTTPException(400, detail="文件格式错误：仅支持 .json 或 .zip 文件")
    raw = await file.read()
    is_zip = file.filename.lower().endswith(".zip")
    if is_zip:
        import zipfile, io
        try:
            with zipfile.ZipFile(io.BytesIO(raw)) as zf:
                raw = zf.read("presentation.json")
        except Exception:
            raise HTTPException(400, detail="ZIP 解析失败：未找到 presentation.json")
    try:
        pres_dict = json.loads(raw)
    except json.JSONDecodeError:
        raise HTTPException(400, detail="JSON 解析失败：文件内容不是有效的 JSON 格式")
    data = Presentation(**pres_dict)
    if slide_index >= len(data.slides):
        raise HTTPException(400, detail=f"页码越界：共 {len(data.slides)} 页，请求第 {slide_index} 页")

    # Build the real PPTX and read actual positions back
    tmp_path = get_temp_path(".pptx")
    try:
        build_pptx(data, tmp_path)
        prs = PptxPresentation(tmp_path)
        if slide_index >= len(prs.slides):
            raise HTTPException(400, detail=f"生成的 PPTX 只有 {len(prs.slides)} 页")
        pptx_slide = prs.slides[slide_index]
        sw = prs.slide_width
        sh = prs.slide_height

        shapes_info = []
        # Collect image URLs from JSON shapes (in order) — consume via queue as we encounter images in PPTX
        image_url_queue = []
        for s in data.slides[slide_index].shapes:
            if s.type.value == "image" and s.image_content and s.image_content.path:
                fname = os.path.basename(s.image_content.path)
                image_url_queue.append(f"/image/{fname}")
        for z, pptx_shape in enumerate(pptx_slide.shapes):
            parsed = _parse_shape(pptx_shape, z)
            if parsed is None:
                continue

            left = int(getattr(pptx_shape, 'left', 0) or 0)
            top = int(getattr(pptx_shape, 'top', 0) or 0)
            width = int(getattr(pptx_shape, 'width', 0) or 0)
            height = int(getattr(pptx_shape, 'height', 0) or 0)

            # Extract text content from the built shape
            paragraphs_data = []
            text_preview = None
            if pptx_shape.has_text_frame:
                for pi, para in enumerate(pptx_shape.text_frame.paragraphs):
                    runs_data = []
                    for run in para.runs:
                        r_text = run.text[:200]
                        color = None
                        try:
                            if run.font.color and run.font.color.rgb:
                                color = str(run.font.color.rgb)
                        except AttributeError:
                            color = None
                        r_info = {
                            "text": r_text,
                            "font_size": run.font.size,
                            "font_name": run.font.name,
                            "bold": run.font.bold,
                            "italic": run.font.italic,
                            "underline": run.font.underline,
                            "color": color,
                        }
                        runs_data.append(r_info)
                        if not text_preview and r_text.strip():
                            text_preview = r_text[:80]
                    if runs_data:
                        paragraphs_data.append({
                            "alignment": str(para.alignment).split('.')[-1].lower() if para.alignment else "left",
                            "line_spacing": para.line_spacing,
                            "runs": runs_data,
                        })

            # Table data from built shape
            table_data = None
            if pptx_shape.shape_type == 19:  # TABLE
                try:
                    tbl = pptx_shape.table
                    cells_data = []
                    for row in tbl.rows:
                        row_cells = []
                        for cell in row.cells:
                            row_cells.append({
                                "text": cell.text.strip() or "",
                                "background_color": "",
                                "bold": False,
                            })
                        cells_data.append(row_cells)
                    table_data = {
                        "rows": len(tbl.rows),
                        "cols": len(tbl.columns),
                        "cells": cells_data,
                    }
                except Exception:
                    pass

            # Match image URL from JSON queue (consumed in order)
            image_content = None
            if parsed.type.value == "image" and image_url_queue:
                image_content = {"url": image_url_queue.pop(0)}

            shapes_info.append({
                "id": parsed.id,
                "role": parsed.role,
                "type": parsed.type.value,
                "auto_shape_type": parsed.auto_shape_type,
                "left": left,
                "top": top,
                "width": width,
                "height": height,
                "background_color": parsed.background_color,
                "border_color": parsed.border_color,
                "border_width": parsed.border_width,
                "rotation": parsed.rotation,
                "z_order": z,
                "text_preview": text_preview,
                "paragraphs": paragraphs_data,
                "table": table_data,
                "image_content": image_content,
            })

        slide_bg = None
        try:
            slide_bg = _parse_xml_fill(pptx_slide._element)
        except Exception:
            pass

        return {
            "slide_index": slide_index,
            "total_slides": len(data.slides),
            "slide_width": sw,
            "slide_height": sh,
            "preset": data.slides[slide_index].preset,
            "background_color": slide_bg,
            "shapes": shapes_info,
        }
    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass


@router.post("/select-element")
async def select_element(request: Request):
    body = await request.json()
    logger.info("Element selected", **body)
    return {"status": "logged"}
