import os
import yaml
from pathlib import Path
from pptx import Presentation as PptxPresentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
from pptx.oxml.ns import qn
from ppt_render_engine.models.schema import (
    Presentation, Slide, Shape, ShapeType, Paragraph, TableContent,
)
from ppt_render_engine.core.template import get_template_manager, TemplateNotFoundError, LayoutNotFoundError
from ppt_render_engine.core.color_scheme import get_color_scheme_manager
from ppt_render_engine.engine.layout_preset import resolve_preset
from ppt_render_engine.engine.validator import validate_slide
from ppt_render_engine.log_config import app_logger
from ppt_render_engine.temp_util import get_temp_path

logger = app_logger("core.builder")

_cfg_path = Path(__file__).parents[2] / "config.yaml"


def _get_image_dir() -> str:
    try:
        with open(_cfg_path, encoding="utf-8") as f:
            cfg = yaml.safe_load(f) or {}
        return cfg.get("image", {}).get("storage_dir", "./sync/images")
    except Exception:
        return "./sync/images"


AUTO_SHAPE_MAP = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "diamond": MSO_SHAPE.DIAMOND,
    "parallelogram": MSO_SHAPE.PARALLELOGRAM,
    "trapezoid": MSO_SHAPE.TRAPEZOID,
    "pentagon": MSO_SHAPE.PENTAGON,
    "regular_pentagon": MSO_SHAPE.REGULAR_PENTAGON,
    "hexagon": MSO_SHAPE.HEXAGON,
    "heptagon": MSO_SHAPE.HEPTAGON,
    "octagon": MSO_SHAPE.OCTAGON,
    "decagon": MSO_SHAPE.DECAGON,
    "dodecagon": MSO_SHAPE.DODECAGON,
    "cross": MSO_SHAPE.CROSS,
    "right_arrow": MSO_SHAPE.RIGHT_ARROW,
    "left_arrow": MSO_SHAPE.LEFT_ARROW,
    "up_arrow": MSO_SHAPE.UP_ARROW,
    "down_arrow": MSO_SHAPE.DOWN_ARROW,
    "chevron": MSO_SHAPE.CHEVRON,
    "bent_arrow": MSO_SHAPE.BENT_ARROW,
    "bent_up_arrow": MSO_SHAPE.BENT_UP_ARROW,
    "circular_arrow": MSO_SHAPE.CIRCULAR_ARROW,
    "curved_right_arrow": MSO_SHAPE.CURVED_RIGHT_ARROW,
    "curved_left_arrow": MSO_SHAPE.CURVED_LEFT_ARROW,
    "curved_up_arrow": MSO_SHAPE.CURVED_UP_ARROW,
    "curved_down_arrow": MSO_SHAPE.CURVED_DOWN_ARROW,
    "striped_right_arrow": MSO_SHAPE.STRIPED_RIGHT_ARROW,
    "notched_right_arrow": MSO_SHAPE.NOTCHED_RIGHT_ARROW,
    "left_right_arrow": MSO_SHAPE.LEFT_RIGHT_ARROW,
    "up_down_arrow": MSO_SHAPE.UP_DOWN_ARROW,
    "quad_arrow": MSO_SHAPE.QUAD_ARROW,
    "right_bracket": MSO_SHAPE.RIGHT_BRACKET,
    "left_bracket": MSO_SHAPE.LEFT_BRACKET,
    "right_brace": MSO_SHAPE.RIGHT_BRACE,
    "left_brace": MSO_SHAPE.LEFT_BRACE,
    "double_bracket": MSO_SHAPE.DOUBLE_BRACKET,
    "double_brace": MSO_SHAPE.DOUBLE_BRACE,
    "sun": MSO_SHAPE.SUN,
    "moon": MSO_SHAPE.MOON,
    "cloud": MSO_SHAPE.CLOUD,
    "heart": MSO_SHAPE.HEART,
    "lightning_bolt": MSO_SHAPE.LIGHTNING_BOLT,
    "star_5_point": MSO_SHAPE.STAR_5_POINT,
    "star_6_point": MSO_SHAPE.STAR_6_POINT,
    "star_7_point": MSO_SHAPE.STAR_7_POINT,
    "star_10_point": MSO_SHAPE.STAR_10_POINT,
    "star_12_point": MSO_SHAPE.STAR_12_POINT,
    "star_16_point": MSO_SHAPE.STAR_16_POINT,
    "star_24_point": MSO_SHAPE.STAR_24_POINT,
    "star_32_point": MSO_SHAPE.STAR_32_POINT,
    "no_symbol": MSO_SHAPE.NO_SYMBOL,
    "pie": MSO_SHAPE.PIE,
    "pie_wedge": MSO_SHAPE.PIE_WEDGE,
    "block_arc": MSO_SHAPE.BLOCK_ARC,
    "donut": MSO_SHAPE.DONUT,
    "bevel": MSO_SHAPE.BEVEL,
    "cube": MSO_SHAPE.CUBE,
    "can": MSO_SHAPE.CAN,
    "folded_corner": MSO_SHAPE.FOLDED_CORNER,
    "corner_tabs": MSO_SHAPE.CORNER_TABS,
    "plaque_tabs": MSO_SHAPE.PLAQUE_TABS,
    "horizontal_scroll": MSO_SHAPE.HORIZONTAL_SCROLL,
    "vertical_scroll": MSO_SHAPE.VERTICAL_SCROLL,
    "wave": MSO_SHAPE.WAVE,
    "double_wave": MSO_SHAPE.DOUBLE_WAVE,
    "tear": MSO_SHAPE.TEAR,
    "balloon": MSO_SHAPE.BALLOON,
    "funnel": MSO_SHAPE.FUNNEL,
    "gear_6": MSO_SHAPE.GEAR_6,
    "gear_9": MSO_SHAPE.GEAR_9,
    "flowchart_process": MSO_SHAPE.FLOWCHART_PROCESS,
    "flowchart_decision": MSO_SHAPE.FLOWCHART_DECISION,
    "flowchart_document": MSO_SHAPE.FLOWCHART_DOCUMENT,
    "flowchart_terminator": MSO_SHAPE.FLOWCHART_TERMINATOR,
    "flowchart_data": MSO_SHAPE.FLOWCHART_DATA,
    "action_button_custom": MSO_SHAPE.ACTION_BUTTON_CUSTOM,
    "math_plus": MSO_SHAPE.MATH_PLUS,
    "math_minus": MSO_SHAPE.MATH_MINUS,
    "math_multiply": MSO_SHAPE.MATH_MULTIPLY,
    "math_divide": MSO_SHAPE.MATH_DIVIDE,
    "math_equal": MSO_SHAPE.MATH_EQUAL,
    "math_not_equal": MSO_SHAPE.MATH_NOT_EQUAL,
}


def build_pptx(presentation: Presentation, output_path: str) -> None:
    tm = get_template_manager()
    template_available = presentation.template_name and tm.has_template(presentation.template_name)
    if template_available:
        logger.info("Using pre-loaded template", template=presentation.template_name)
        prs = PptxPresentation(tm.filepath_for(presentation.template_name))
    else:
        if presentation.template_name:
            logger.warning("Template not loaded, falling back to blank", template=presentation.template_name)
        prs = PptxPresentation()
    _set_slide_size(prs, presentation.slide_width, presentation.slide_height)
    csm = get_color_scheme_manager()
    scheme_name = presentation.color_scheme or csm.get_default_name()
    csm.apply_to_presentation(prs, scheme_name)
    original_slide_count = len(prs.slides)
    for i, slide_data in enumerate(presentation.slides):
        slide_data = resolve_preset(slide_data)
        if i < original_slide_count:
            pptx_slide = prs.slides[i]
            layout = pptx_slide.slide_layout
        else:
            layout = _resolve_layout(prs, slide_data, presentation.template_name)
            pptx_slide = prs.slides.add_slide(layout)
        slide_data = _resolve_placeholder_coords(slide_data, layout)
        slide_data = _apply_fallback_coords(slide_data, presentation)
        warnings = validate_slide(slide_data, slide_index=i)
        for w in warnings:
            logger.warning("Layout validation", slide=i, issue=w)
        layout_name = slide_data.layout_id or (slide_data.preset or "blank")
        logger.debug("Processing slide", index=i, layout=layout_name)
        _build_slide(pptx_slide, slide_data, csm, scheme_name)
    # remove excess template slides
    for _ in range(original_slide_count - len(presentation.slides)):
        sldIdLst = prs.slides._sldIdLst
        rId = sldIdLst[-1].get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldIdLst[-1])
    logger.info("Saved PPTX", path=output_path)
    tmp_path = output_path + ".tmp"
    prs.save(tmp_path)
    os.replace(tmp_path, output_path)

def _resolve_layout(prs: PptxPresentation, slide_data: Slide, template_name: str | None) -> object:
    if template_name and slide_data.layout_id:
        tm = get_template_manager()
        try:
            idx = tm.get_layout_index(template_name, slide_data.layout_id)
            return prs.slide_layouts[idx]
        except (TemplateNotFoundError, LayoutNotFoundError):
            pass
    if slide_data.layout_id:
        name = slide_data.layout_id.strip().lower()
        for i, layout in enumerate(prs.slide_layouts):
            if layout.name.strip().lower() == name:
                return prs.slide_layouts[i]
    return prs.slide_layouts[6]


def _resolve_placeholder_coords(slide_data: Slide, layout) -> Slide:
    """Fill left/top/width/height from template layout placeholders for shapes with placeholder/role refs.
    Only fills if the shape has no explicit coordinates (all zero)."""
    resolved = slide_data.model_copy(deep=True)
    ph_positions = _get_layout_placeholder_positions(layout)
    for shape in resolved.shapes:
        ref = shape.placeholder or shape.role
        if ref and ref in ph_positions:
            pos = ph_positions[ref]
            if not (shape.left or shape.top or shape.width or shape.height):
                shape.left = pos["left"]
                shape.top = pos["top"]
                shape.width = pos["width"]
                shape.height = pos["height"]
    return resolved


def _get_layout_placeholder_positions(layout) -> dict[str, dict[str, float]]:
    positions = {}
    try:
        for ph in layout.placeholders:
            key = str(ph.placeholder_format.idx)
            positions[key] = {
                "left": ph.left,
                "top": ph.top,
                "width": ph.width,
                "height": ph.height,
            }
            positions[ph.name.lower()] = positions[key]
    except Exception:
        pass
    return positions


def _apply_fallback_coords(slide_data: Slide, presentation: Presentation) -> Slide:
    resolved = slide_data.model_copy(deep=True)
    margin = presentation.slide_width * 0.04 if presentation.slide_width else 500000
    y = margin
    for shape in resolved.shapes:
        if not (shape.placeholder or shape.role):
            continue
        if shape.width and shape.height:
            y += shape.height + margin
            continue
        shape.left = int(margin)
        shape.top = int(y)
        shape.width = int(presentation.slide_width - margin * 2)
        shape.height = int(presentation.slide_height * 0.08)
        y += shape.height + margin
    return resolved


def _set_slide_size(prs, width: float, height: float):
    if width and height:
        prs.slide_width = int(width)
        prs.slide_height = int(height)


def _ensure_grp_sp_pr(pptx_slide):
    """Ensure spTree/grpSpPr has a:xfrm (required by some PowerPoint versions)."""
    spTree = pptx_slide._element.find(qn('p:cSld') + '/' + qn('p:spTree'))
    if spTree is None:
        return
    grpSpPr = spTree.find(qn('p:grpSpPr'))
    if grpSpPr is None:
        grpSpPr = etree.SubElement(spTree, qn('p:grpSpPr'))
    if grpSpPr.find(qn('a:xfrm')) is None:
        xfrm = etree.SubElement(grpSpPr, qn('a:xfrm'))
        etree.SubElement(xfrm, qn('a:off'), x="0", y="0")
        etree.SubElement(xfrm, qn('a:ext'), cx="0", cy="0")
        etree.SubElement(xfrm, qn('a:chOff'), x="0", y="0")
        etree.SubElement(xfrm, qn('a:chExt'), cx="0", cy="0")


def _build_slide(pptx_slide, slide_data: Slide, csm, scheme_name: str):
    _ensure_grp_sp_pr(pptx_slide)
    if slide_data.background_color:
        _set_slide_background(pptx_slide, slide_data.background_color)
    ph_map = _build_placeholder_map(pptx_slide)
    for shape_data in sorted(slide_data.shapes, key=lambda s: s.z_order or 0):
        ref = shape_data.placeholder or shape_data.role
        if ref and ref.lower() in ph_map:
            target = ph_map[ref.lower()]
            # For table/shape types that can't reuse a text placeholder,
            # inherit the placeholder's position and create a new element.
            if shape_data.type in (ShapeType.TABLE, ShapeType.SHAPE):
                _build_shape(pptx_slide, shape_data, csm, scheme_name)
                continue
            if shape_data.text_content and target.has_text_frame:
                tf = target.text_frame
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.NONE
                bodyPr = tf._txBody.find(qn('a:bodyPr'))
                if bodyPr is not None:
                    bodyPr.set('vertOverflow', 'clip')
                existing = list(tf.paragraphs)
                for i, para_data in enumerate(shape_data.text_content):
                    if i < len(existing):
                        p = existing[i]
                    else:
                        p = tf.add_paragraph()
                    for run in list(p.runs):
                        run._r.getparent().remove(run._r)
                    _build_paragraph(p, para_data, csm, scheme_name)
                for j in range(len(shape_data.text_content), len(existing)):
                    existing[j]._p.getparent().remove(existing[j]._p)
            if shape_data.background_color:
                spPr = target._element.find(qn('p:spPr'))
                if spPr is not None:
                    _apply_xml_solid_fill(spPr, shape_data.background_color, csm, scheme_name)
            if shape_data.border_color or shape_data.border_width:
                try:
                    _apply_line(target.line, shape_data, csm, scheme_name)
                except Exception:
                    pass
            continue
        _build_shape(pptx_slide, shape_data, csm, scheme_name)


def _build_placeholder_map(pptx_slide) -> dict[str, object]:
    ph_map = {}
    for shape in pptx_slide.shapes:
        if shape.is_placeholder:
            idx = str(shape.placeholder_format.idx)
            ph_map[idx] = shape
            ph_map[shape.name.lower()] = shape
    return ph_map


def _build_shape(pptx_slide, shape_data: Shape, csm, scheme_name: str):
    left = int(shape_data.left)
    top = int(shape_data.top)
    width = max(int(shape_data.width), 1)
    height = max(int(shape_data.height), 1)

    if shape_data.type == ShapeType.TABLE and shape_data.table_content:
        table_shape = _build_table(pptx_slide, shape_data.table_content, left, top, width, height, csm, scheme_name)
        if table_shape is not None and shape_data.id:
            table_shape.name = shape_data.id
        return
    if shape_data.type == ShapeType.IMAGE and shape_data.image_content:
        _build_image(pptx_slide, shape_data.image_content, left, top, width, height, shape_data.id)
        return
    if shape_data.type == ShapeType.SHAPE and shape_data.auto_shape_type:
        mso = AUTO_SHAPE_MAP.get(shape_data.auto_shape_type)
        if mso:
            shape = pptx_slide.shapes.add_shape(mso, left, top, width, height)
            shape.name = shape_data.id or f"shape_{shape_data.z_order}"
            _apply_shape_style(shape, shape_data, csm, scheme_name)
            if shape_data.text_content:
                _build_text_frame(shape.text_frame, shape_data.text_content, csm, scheme_name)
            else:
                _remove_empty_txBody(shape)
            return

    txBox = pptx_slide.shapes.add_textbox(left, top, width, height)
    txBox.name = shape_data.id or f"shape_{shape_data.z_order}"
    if shape_data.rotation:
        txBox.rotation = shape_data.rotation
    if shape_data.text_content:
        _build_text_frame(txBox.text_frame, shape_data.text_content, csm, scheme_name)
    if shape_data.background_color:
        spPr = txBox._element.find(qn('p:spPr'))
        if spPr is None:
            spPr = etree.SubElement(txBox._element, qn('p:spPr'))
        _apply_xml_solid_fill(spPr, shape_data.background_color, csm, scheme_name)
    if shape_data.border_color or shape_data.border_width:
        _apply_line(txBox.line, shape_data, csm, scheme_name)


_FILL_TAGS = (qn('a:solidFill'), qn('a:noFill'), qn('a:gradFill'),
              qn('a:blipFill'), qn('a:pattFill'), qn('a:grpFill'))


def _apply_xml_solid_fill(parent_elem, hex_color, csm, scheme_name):
    """Remove existing fills and add a solidFill using schemeClr if it matches the scheme."""
    for child in list(parent_elem):
        if child.tag in _FILL_TAGS:
            parent_elem.remove(child)
    sf = etree.SubElement(parent_elem, qn('a:solidFill'))
    color_name = csm.get_color_name(hex_color, scheme_name)
    if color_name:
        etree.SubElement(sf, qn('a:schemeClr'), val=color_name)
    else:
        etree.SubElement(sf, qn('a:srgbClr'), val=hex_color.upper().lstrip("#"))


def _apply_shape_style(shape, shape_data: Shape, csm, scheme_name: str):
    if shape_data.rotation:
        shape.rotation = shape_data.rotation
    if shape_data.background_color:
        spPr = shape._element.find(qn('p:spPr'))
        if spPr is not None:
            _apply_xml_solid_fill(spPr, shape_data.background_color, csm, scheme_name)
    if shape_data.border_color or shape_data.border_width:
        _apply_line(shape.line, shape_data, csm, scheme_name)


def _ensure_ln(shape) -> object:
    """Ensure an <a:ln> element exists on the shape and return it."""
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        spPr = etree.SubElement(shape._element, qn('p:spPr'))
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        ln = etree.SubElement(spPr, qn('a:ln'))
    return ln


def _apply_line(line, shape_data: Shape, csm, scheme_name: str):
    if shape_data.border_color:
        color_name = csm.get_color_name(shape_data.border_color, scheme_name)
        ln = _ensure_ln(line._parent)
        for child in list(ln):
            ln.remove(child)
        sf = etree.SubElement(ln, qn('a:solidFill'))
        if color_name:
            etree.SubElement(sf, qn('a:schemeClr'), val=color_name)
        else:
            hex_val = shape_data.border_color.upper().lstrip("#")
            etree.SubElement(sf, qn('a:srgbClr'), val=hex_val)
    if shape_data.border_width is not None:
        try:
            line.width = Pt(shape_data.border_width / 12700)
        except Exception:
            pass


def _remove_empty_txBody(shape):
    txBody = shape._element.find(qn('p:txBody'))
    if txBody is not None:
        shape._element.remove(txBody)


def _build_image(pptx_slide, image_data, left, top, width, height, shape_id=None):
    import base64
    import urllib.request
    img_path = None
    if image_data.path:
        if '/' not in image_data.path and '\\' not in image_data.path:
            candidate = os.path.join(_get_image_dir(), image_data.path)
            if os.path.exists(candidate):
                img_path = candidate
        elif os.path.exists(image_data.path):
            img_path = image_data.path
    if img_path is None and image_data.data:
        raw = base64.b64decode(image_data.data)
        img_path = get_temp_path(".png")
        with open(img_path, "wb") as f:
            f.write(raw)
    elif img_path is None and image_data.url:
        try:
            img_path = get_temp_path(".png")
            urllib.request.urlretrieve(image_data.url, img_path)
        except Exception:
            img_path = None
    cleanup = img_path is not None and (image_data.data or image_data.url)
    if img_path:
        pic = pptx_slide.shapes.add_picture(img_path, left, top, width, height)
        if shape_id:
            pic.name = shape_id
        if cleanup:
            os.unlink(img_path)
        return
    # Fallback: add placeholder rectangle so block stays visible
    shape = pptx_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if shape_id:
        shape.name = shape_id
    _remove_empty_txBody(shape)
    shape.fill.background()
    try:
        shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        shape.line.width = Pt(0.5)
    except Exception:
        pass


def _build_text_frame(tf, paragraphs: list[Paragraph], csm, scheme_name: str):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('vertOverflow', 'clip')
    for i, para_data in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        _build_paragraph(p, para_data, csm, scheme_name)


def _build_paragraph(p, para_data: Paragraph, csm, scheme_name: str):
    align_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }
    p.alignment = align_map.get(para_data.alignment)
    if para_data.line_spacing is not None:
        p.line_spacing = Pt(para_data.line_spacing / 12700)
    if para_data.space_before is not None:
        p.space_before = int(para_data.space_before)
    if para_data.space_after is not None:
        p.space_after = int(para_data.space_after)
    for run_data in para_data.runs:
        run = p.add_run()
        run.text = run_data.text
        if run_data.font_size is not None:
            run.font.size = run_data.font_size
        if run_data.font_name is not None:
            run.font.name = run_data.font_name
        if run_data.bold:
            run.font.bold = True
        if run_data.italic:
            run.font.italic = True
        if run_data.underline:
            run.font.underline = True
        if run_data.color is not None:
            color_name = csm.get_color_name(run_data.color, scheme_name)
            if color_name:
                rPr = run._r.find(qn('a:rPr'))
                if rPr is None:
                    rPr = etree.SubElement(run._r, qn('a:rPr'))
                if rPr.find(qn('a:solidFill')) is None:
                    fill = etree.SubElement(rPr, qn('a:solidFill'))
                    etree.SubElement(fill, qn('a:schemeClr'), val=color_name)
            else:
                run.font.color.rgb = RGBColor.from_string(run_data.color.lstrip("#"))
        else:
            rPr = run._r.find(qn('a:rPr'))
            if rPr is None:
                rPr = etree.SubElement(run._r, qn('a:rPr'))
            if rPr.find(qn('a:solidFill')) is None:
                fill = etree.SubElement(rPr, qn('a:solidFill'))
                etree.SubElement(fill, qn('a:schemeClr'), val="dk1")


def _build_table(pptx_slide, table_data: TableContent, left, top, width, height, csm, scheme_name: str):
    rows, cols = table_data.rows, table_data.cols
    table_shape = pptx_slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    for r_idx, row_data in enumerate(table_data.cells):
        for c_idx, cell_data in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = cell_data.text
            if cell_data.background_color:
                color_name = csm.get_color_name(cell_data.background_color, scheme_name)
                tc = cell._tc
                tcPr = tc.find(qn('a:tcPr'))
                if tcPr is None:
                    tcPr = etree.SubElement(tc, qn('a:tcPr'))
                for child in list(tcPr):
                    if child.tag in _FILL_TAGS:
                        tcPr.remove(child)
                sf = etree.SubElement(tcPr, qn('a:solidFill'))
                if color_name:
                    etree.SubElement(sf, qn('a:schemeClr'), val=color_name)
                else:
                    hex_val = cell_data.background_color.upper().lstrip("#")
                    etree.SubElement(sf, qn('a:srgbClr'), val=hex_val)
    return table_shape


def _set_slide_background(pptx_slide, color: str):
    bg = pptx_slide.background
    fill = bg.fill
    fill.solid()
    try:
        fill.fore_color.rgb = RGBColor.from_string(color.lstrip("#"))
    except Exception:
        pass


def build_json_from_pptx(filepath: str) -> dict:
    from ppt_render_engine.core.pptx_parser import parse_pptx
    pres = parse_pptx(filepath)
    return pres.model_dump(exclude_none=True)
