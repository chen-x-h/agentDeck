import os
import yaml
from pathlib import Path
from pptx import Presentation as PptxPresentation
from pptx.oxml.ns import qn
from ppt_render_engine.log_config import app_logger
from ppt_render_engine.temp_util import get_temp_path

logger = app_logger("core.template")
_cfg_path = Path(__file__).parents[3] / "config.yaml"


class TemplateNotFoundError(KeyError):
    pass


class LayoutNotFoundError(KeyError):
    pass


class TemplateManager:
    def __init__(self):
        self._templates: dict[str, PptxPresentation] = {}
        self._layout_map: dict[str, dict[str, int]] = {}

    def load(self, name: str, filepath: str) -> None:
        if not os.path.isfile(filepath):
            raise FileNotFoundError(f"Template file not found: {filepath}")
        logger.info("Loading template", name=name, path=filepath)
        prs = PptxPresentation(filepath)
        self._templates[name] = prs
        layouts = {}
        for i, layout in enumerate(prs.slide_layouts):
            key = _normalize_layout_name(layout.name)
            layouts[key] = i
            logger.debug("  layout", index=i, name=layout.name)
        self._layout_map[name] = layouts
        logger.info("Template loaded", name=name, layouts=len(layouts))

    def unload(self, name: str) -> None:
        self._templates.pop(name, None)
        self._layout_map.pop(name, None)
        logger.info("Template unloaded", name=name)

    def list_templates(self) -> list[str]:
        return list(self._templates.keys())

    def list_layouts(self, name: str) -> list[dict]:
        if name not in self._templates:
            raise TemplateNotFoundError(name)
        prs = self._templates[name]
        results = []
        for i, layout in enumerate(prs.slide_layouts):
            placeholders = []
            for ph in layout.placeholders:
                ptype = str(ph.placeholder_format.type)
                if ptype in ("DT", "FTR", "SLD_NUM"):
                    continue
                placeholders.append({
                    "idx": ph.placeholder_format.idx,
                    "name": ph.name,
                    "type": str(ph.placeholder_format.type),
                    "left": ph.left,
                    "top": ph.top,
                    "width": ph.width,
                    "height": ph.height,
                })
            results.append({"index": i, "name": layout.name, "placeholders": placeholders})
        return results

    def get_presentation(self, name: str) -> PptxPresentation:
        if name not in self._templates:
            raise TemplateNotFoundError(name)
        return self._templates[name]

    def get_layout_index(self, template_name: str, layout_name_or_idx: str) -> int:
        if template_name not in self._layout_map:
            raise TemplateNotFoundError(template_name)
        # try as index first
        try:
            idx = int(layout_name_or_idx)
            if 0 <= idx < len(self._layout_map[template_name]):
                return idx
        except (ValueError, TypeError):
            pass
        # try exact normalized match
        key = _normalize_layout_name(layout_name_or_idx)
        layout_map = self._layout_map[template_name]
        if key in layout_map:
            return layout_map[key]
        # try bilingual alias
        alias = _BILINGUAL_LAYOUT.get(key)
        if alias and alias in layout_map:
            return layout_map[alias]
        # try fuzzy: match if key is contained in any layout key
        for k, v in layout_map.items():
            if key in k or k in key:
                return v
        available = list(layout_map.keys())
        raise LayoutNotFoundError(
            f"Layout '{layout_name_or_idx}' not found in template '{template_name}'. Available: {available}"
        )

    def has_template(self, name: str) -> bool:
        return name in self._templates

    @staticmethod
    def get_storage_dir() -> str:
        try:
            with open(_cfg_path, encoding="utf-8") as f:
                cfg = yaml.safe_load(f) or {}
            return cfg.get("template", {}).get("storage_dir", "./storage/templates")
        except Exception:
            return "./storage/templates"

    def filepath_for(self, name: str) -> str:
        d = self.get_storage_dir()
        os.makedirs(d, exist_ok=True)
        return os.path.abspath(os.path.join(d, f"{name}.pptx"))

    def load_from_storage(self, name: str) -> bool:
        path = self.filepath_for(name)
        if not os.path.isfile(path):
            return False
        self.load(name, path)
        return True

    def save_to_storage(self, name: str, data: bytes) -> str:
        path = self.filepath_for(name)
        os.makedirs(os.path.dirname(path), exist_ok=True)
        tmp = get_temp_path(".pptx")
        try:
            with open(tmp, "wb") as f:
                f.write(data)
            prs = PptxPresentation(tmp)
            sldIdLst = prs.slides._sldIdLst
            for sld in list(sldIdLst):
                rId = sld.get(qn("r:id"))
                if rId:
                    prs.part.drop_rel(rId)
                sldIdLst.remove(sld)
            prs.save(path)
        finally:
            if os.path.isfile(tmp):
                os.unlink(tmp)
        logger.info("Template saved to storage (slides stripped)", name=name, path=path)
        return path

    def delete_from_storage(self, name: str) -> None:
        path = self.filepath_for(name)
        if os.path.isfile(path):
            os.unlink(path)
            logger.info("Template file deleted", name=name, path=path)

    def preload_all(self) -> list[str]:
        d = self.get_storage_dir()
        if not os.path.isdir(d):
            logger.info("Storage dir not found, skipping preload", path=d)
            return []
        loaded = []
        for fname in os.listdir(d):
            if fname.endswith(".pptx"):
                name = fname[:-5]
                try:
                    self.load_from_storage(name)
                    loaded.append(name)
                    logger.info("Preloaded template", name=name)
                except Exception as e:
                    logger.error("Failed to preload template", name=name, error=str(e))
        logger.info("Templates preloaded", count=len(loaded))
        return loaded


_templates: TemplateManager | None = None


def get_template_manager() -> TemplateManager:
    global _templates
    if _templates is None:
        _templates = TemplateManager()
        logger.debug("TemplateManager singleton created")
    return _templates


def _normalize_layout_name(name: str) -> str:
    return name.strip().lower().replace(" ", "_").replace("\t", "_")


_BILINGUAL_LAYOUT: dict[str, str] = {
    "title_slide": "标题幻灯片",
    "title_and_content": "标题和内容",
    "section_header": "节标题",
    "two_content": "两栏内容",
    "comparison": "比较",
    "title_only": "仅标题",
    "blank": "空白",
    "content_with_caption": "内容与标题",
    "picture_with_caption": "图片与标题",
    "title_and_vertical_text": "标题和竖排文字",
    "vertical_title_and_text": "垂直排列标题与\n文本",
}
