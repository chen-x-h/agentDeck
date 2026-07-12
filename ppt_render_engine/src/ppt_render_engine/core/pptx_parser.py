import os
from pptx import Presentation as PptxPresentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
from ppt_render_engine.models.schema import (
    Presentation, Slide, Shape, ShapeType, Paragraph, TextRun,
    ImageContent, TableContent, CellContent,
)
from pptx.oxml.ns import qn
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from ppt_render_engine.log_config import app_logger

logger = app_logger("core.parser")


def _extract_theme_colors(prs: PptxPresentation) -> dict[str, str]:
    """Extract clrScheme hex values from theme: {name: hex}."""
    try:
        sm = prs.slide_masters[0]
        for rel in sm.part.rels.values():
            if rel.reltype == RT.THEME:
                theme = etree.fromstring(rel.target_part.blob)
                clrScheme = theme.find('.//' + qn('a:clrScheme'))
                if clrScheme is not None:
                    colors = {}
                    for child in clrScheme:
                        tag = child.tag.split('}')[-1]
                        srgb = child.find(qn('a:srgbClr'))
                        if srgb is not None:
                            colors[tag] = srgb.get('val')
                            continue
                        sysClr = child.find(qn('a:sysClr'))
                        if sysClr is not None:
                            colors[tag] = sysClr.get('lastClr')
                    return colors
    except Exception:
        pass
    return {}


def parse_pptx(filepath: str) -> Presentation:
    logger.info("Opening PPTX file", path=filepath)
    prs = PptxPresentation(filepath)
    theme_colors = _extract_theme_colors(prs)
    sw = prs.slide_width
    sh = prs.slide_height
    slides = []
    for idx, pptx_slide in enumerate(prs.slides):
        logger.debug("Parsing slide", index=idx)
        slide = _parse_slide(pptx_slide, idx, sw, sh, theme_colors)
        slides.append(slide)
    result = Presentation(
        title=prs.core_properties.title,
        slide_width=prs.slide_width,
        slide_height=prs.slide_height,
        slides=slides,
    )
    logger.info("PPTX parsed", slides=len(slides), title=result.title)
    return result


def _parse_slide(pptx_slide, index: int, slide_w: float, slide_h: float,
                 theme_colors: dict[str, str] | None = None) -> Slide:
    slide = Slide(
        id=index,
        width=slide_w,
        height=slide_h,
        layout_id=pptx_slide.slide_layout.name if pptx_slide.slide_layout else None,
    )
    if pptx_slide.background.fill.type is not None:
        slide.background_color = _parse_fill_color(pptx_slide.background)
    for z, shape in enumerate(pptx_slide.shapes):
        parsed = _parse_shape(shape, z, theme_colors)
        if parsed is not None:
            slide.shapes.append(parsed)
    logger.debug("Slide parsed", index=index, shapes=len(slide.shapes))
    return slide


def _parse_shape(shape, z_order: int, theme_colors: dict[str, str] | None = None) -> Shape | None:
    if shape.shape_type is None:
        return None
    st = _map_shape_type(shape.shape_type)
    if st is None:
        return None
    # Use shape name as id (builder stores semantic id there); fall back to shape_id
    sid = getattr(shape, 'name', None) or str(getattr(shape, 'shape_id', z_order))
    # Capture placeholder reference so builder can fill layout placeholders correctly
    placeholder_ref = None
    nvPr = shape._element.find('.//' + qn('p:nvPr'))
    if nvPr is not None:
        ph = nvPr.find(qn('p:ph'))
        if ph is not None:
            placeholder_ref = ph.get('idx', '0')
    s = Shape(
        placeholder=placeholder_ref,
        id=sid,
        type=st,
        left=shape.left,
        top=shape.top,
        width=shape.width,
        height=shape.height,
        z_order=z_order,
        rotation=getattr(shape, 'rotation', 0),
    )
    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        s.auto_shape_type = _resolve_auto_shape_name(shape)
    if shape.has_text_frame:
        text_content = _parse_text_frame(shape.text_frame)
        if st == ShapeType.SHAPE and text_content:
            has_content = any(
                any(run.text.strip() for run in p.runs)
                for p in text_content
            )
            if has_content:
                s.text_content = text_content
        else:
            s.text_content = text_content
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        s.image_content = _parse_image(shape)
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        try:
            s.table_content = _parse_table(shape.table)
        except Exception:
            pass
    # Parse fill from XML (more reliable than python-pptx fill API)
    bg = _parse_xml_fill(shape._element, theme_colors)
    if bg:
        s.background_color = bg
    return s


def _resolve_auto_shape_name(shape) -> str | None:
    from pptx.enum.shapes import MSO_SHAPE
    A = MSO_SHAPE
    REVERSE_AUTO_MAP = {
        A.RECTANGLE: "rectangle",
        A.ROUNDED_RECTANGLE: "rounded_rectangle",
        A.OVAL: "oval",
        A.DIAMOND: "diamond",
        A.PARALLELOGRAM: "parallelogram",
        A.TRAPEZOID: "trapezoid",
        A.PENTAGON: "pentagon",
        A.REGULAR_PENTAGON: "regular_pentagon",
        A.HEXAGON: "hexagon",
        A.HEPTAGON: "heptagon",
        A.OCTAGON: "octagon",
        A.DECAGON: "decagon",
        A.DODECAGON: "dodecagon",
        A.CROSS: "cross",
        A.RIGHT_ARROW: "right_arrow",
        A.LEFT_ARROW: "left_arrow",
        A.UP_ARROW: "up_arrow",
        A.DOWN_ARROW: "down_arrow",
        A.CHEVRON: "chevron",
        A.RIGHT_BRACKET: "right_bracket",
        A.LEFT_BRACKET: "left_bracket",
        A.RIGHT_BRACE: "right_brace",
        A.LEFT_BRACE: "left_brace",
        A.DOUBLE_BRACKET: "double_bracket",
        A.DOUBLE_BRACE: "double_brace",
        A.SUN: "sun",
        A.MOON: "moon",
        A.CLOUD: "cloud",
        A.HEART: "heart",
        A.LIGHTNING_BOLT: "lightning_bolt",
        A.STAR_5_POINT: "star_5_point",
        A.STAR_6_POINT: "star_6_point",
        A.STAR_7_POINT: "star_7_point",
        A.STAR_10_POINT: "star_10_point",
        A.STAR_12_POINT: "star_12_point",
        A.STAR_16_POINT: "star_16_point",
        A.STAR_24_POINT: "star_24_point",
        A.STAR_32_POINT: "star_32_point",
        A.NO_SYMBOL: "no_symbol",
        A.PIE: "pie",
        A.DONUT: "donut",
        A.BEVEL: "bevel",
        A.CUBE: "cube",
        A.CAN: "can",
        A.FOLDED_CORNER: "folded_corner",
        A.WAVE: "wave",
        A.DOUBLE_WAVE: "double_wave",
        A.TEAR: "tear",
        A.BALLOON: "balloon",
        A.FUNNEL: "funnel",
        A.GEAR_6: "gear_6",
        A.GEAR_9: "gear_9",
        A.FLOWCHART_PROCESS: "flowchart_process",
        A.FLOWCHART_DECISION: "flowchart_decision",
        A.FLOWCHART_DOCUMENT: "flowchart_document",
        A.FLOWCHART_TERMINATOR: "flowchart_terminator",
        A.FLOWCHART_DATA: "flowchart_data",
        A.ACTION_BUTTON_CUSTOM: "action_button_custom",
        A.MATH_PLUS: "math_plus",
        A.MATH_MINUS: "math_minus",
        A.MATH_MULTIPLY: "math_multiply",
        A.MATH_DIVIDE: "math_divide",
        A.MATH_EQUAL: "math_equal",
        A.MATH_NOT_EQUAL: "math_not_equal",
    }
    try:
        return REVERSE_AUTO_MAP.get(shape.auto_shape_type)
    except Exception:
        return None


def _map_shape_type(st) -> ShapeType:
    mapping = {
        1: ShapeType.SHAPE,    # AUTO_SHAPE
        6: ShapeType.SHAPE,    # GROUP
        13: ShapeType.IMAGE,   # PICTURE
        14: ShapeType.TEXTBOX, # PLACEHOLDER
        17: ShapeType.TEXTBOX, # TEXT_BOX
        19: ShapeType.TABLE,   # TABLE
    }
    return mapping.get(st, ShapeType.SHAPE)


def _parse_text_frame(tf) -> list[Paragraph]:
    paragraphs = []
    for para in tf.paragraphs:
        runs = []
        for run in para.runs:
            color = None
            try:
                if run.font.color and run.font.color.rgb:
                    color = str(run.font.color.rgb)
            except AttributeError:
                color = None
            tr = TextRun(
                text=run.text,
                font_size=run.font.size,
                font_name=run.font.name,
                bold=run.font.bold,
                italic=run.font.italic,
                underline=run.font.underline,
                color=color,
            )
            runs.append(tr)
        align_map = {
            PP_ALIGN.LEFT: "left",
            PP_ALIGN.CENTER: "center",
            PP_ALIGN.RIGHT: "right",
            PP_ALIGN.JUSTIFY: "justify",
        }
        paragraphs.append(Paragraph(
            runs=runs,
            alignment=align_map.get(para.alignment, "left"),
            line_spacing=para.line_spacing,
            space_before=para.space_before,
            space_after=para.space_after,
            indent_level=para.level,
        ))
    return paragraphs


def _parse_image(shape) -> ImageContent:
    img = ImageContent(width=shape.width, height=shape.height)
    try:
        # Extract the image filename from the relationship
        blip = shape._element.find('.//' + qn('a:blip'))
        if blip is not None:
            embed = blip.get(qn('r:embed')) or blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
            if embed:
                rel = shape.part.rels.get(embed)
                if rel and rel.target_part:
                    img.path = os.path.basename(rel.target_part.partname)
        # Fall back to image blob-based naming
        if not img.path:
            ext_map = {"image/png": ".png", "image/jpeg": ".jpg", "image/gif": ".gif"}
            img_obj = shape.image
            ext = ext_map.get(img_obj.content_type, ".png")
            img.path = f"image_{hash(img_obj.blob) & 0xFFFFFFFF:08x}{ext}"
    except Exception:
        pass
    return img


def _parse_table(table) -> TableContent:
    rows, cols = len(table.rows), len(table.columns)
    cells = []
    for row in table.rows:
        row_cells = []
        for cell in row.cells:
            row_cells.append(CellContent(
                text=cell.text,
                colspan=int(cell._tc.get('gridSpan', 1)),
                rowspan=int(cell._tc.get('rowSpan', 1)),
            ))
        cells.append(row_cells)
    return TableContent(rows=rows, cols=cols, cells=cells)


def _parse_xml_fill(elem, theme_colors: dict[str, str] | None = None) -> str | None:
    spPr = elem.find(qn('p:spPr'))
    if spPr is None:
        spPr = elem.find(qn('a:spPr'))
    if spPr is None:
        return None
    sf = spPr.find(qn('a:solidFill'))
    if sf is None:
        return None
    srgb = sf.find(qn('a:srgbClr'))
    if srgb is not None:
        return srgb.get('val')
    scheme = sf.find(qn('a:schemeClr'))
    if scheme is not None:
        name = scheme.get('val')
        # Resolve scheme color to actual hex from theme
        if theme_colors and name in theme_colors:
            return theme_colors[name]
        # Fallback: known common scheme colors
        fallback = {
            "dk1": "000000", "lt1": "FFFFFF",
            "dk2": "44546A", "lt2": "E7E6E6",
            "accent1": "4472C4", "accent2": "ED7D31",
            "accent3": "A5A5A5", "accent4": "FFC000",
            "accent5": "5B9BD5", "accent6": "70AD47",
            "hlink": "0563C1", "folHlink": "954F72",
        }
        return fallback.get(name)
    return None


def _parse_fill_color(fill) -> str | None:
    try:
        if fill.type is not None:
            return str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def _parse_line_color(line) -> str | None:
    try:
        if line.color and line.color.rgb:
            return str(line.color.rgb)
    except Exception:
        pass
    return None
